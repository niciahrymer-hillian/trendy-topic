"""Generate the Trendy Topic capstone PowerPoint from the slide outline.

Run: python scripts/build_slides.py
Output: docs/Trendy_Topic_Capstone.pptx

The content mirrors docs/slide_outline.md. Re-run to regenerate the .pptx.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --- Brand palette (matches the dark dashboard theme) --------------------------
BG = RGBColor(0x0F, 0x14, 0x19)        # near-black background
PANEL = RGBColor(0x1A, 0x22, 0x30)     # panel/footer band
TEXT = RGBColor(0xE6, 0xED, 0xF3)      # primary text
MUTED = RGBColor(0x93, 0xA1, 0xB1)     # secondary text
ACCENT = RGBColor(0x4A, 0xA8, 0xFF)    # blue accent
ACCENT2 = RGBColor(0x7C, 0x5C, 0xFF)   # purple accent

FONT = "Calibri"

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_bg(slide) -> None:
    rect = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    _fill(rect, BG)
    rect.shadow.inherit = False
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def _accent_bar(slide, color: RGBColor = ACCENT) -> None:
    bar = slide.shapes.add_shape(1, 0, Inches(1.55), Inches(4.4), Pt(5))
    _fill(bar, color)
    bar.shadow.inherit = False


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _set(run, size, color=TEXT, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def _footer(slide, idx: int, total: int) -> None:
    band = slide.shapes.add_shape(1, 0, Inches(7.05), SLIDE_W, Inches(0.45))
    _fill(band, PANEL)
    band.shadow.inherit = False
    box, tf = _textbox(slide, Inches(0.4), Inches(7.06), Inches(12.5), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Trendy Topic — What the World Asks AI"
    _set(r, 10, MUTED)
    box2, tf2 = _textbox(slide, Inches(11.8), Inches(7.06), Inches(1.1), Inches(0.4))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{idx} / {total}"
    _set(r2, 10, MUTED)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    # Accent block
    block = slide.shapes.add_shape(1, 0, Inches(2.7), Inches(0.35), Inches(2.1))
    _fill(block, ACCENT2)
    block.shadow.inherit = False

    box, tf = _textbox(slide, Inches(0.9), Inches(2.55), Inches(11.5), Inches(2.6))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Trendy Topic"
    _set(r, 54, TEXT, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = "What the World Asks AI"
    _set(r2, 30, ACCENT, bold=True)
    p3 = tf.add_paragraph()
    r3 = p3.add_run(); r3.text = "A Global AI Conversation Analytics Platform"
    _set(r3, 18, MUTED)

    box4, tf4 = _textbox(slide, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6))
    p4 = tf4.paragraphs[0]
    r4 = p4.add_run()
    r4.text = "Shocka & Niciah  ·  Capstone  ·  2026"
    _set(r4, 16, TEXT, bold=True)
    return slide


def content_slide(prs, idx, total, kicker, title, lead, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)

    # Kicker (section label)
    kbox, ktf = _textbox(slide, Inches(0.6), Inches(0.45), Inches(11), Inches(0.5))
    kp = ktf.paragraphs[0]
    kr = kp.add_run(); kr.text = kicker.upper()
    _set(kr, 14, ACCENT, bold=True)

    # Title
    tbox, ttf = _textbox(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.9))
    tp = ttf.paragraphs[0]
    tr = tp.add_run(); tr.text = title
    _set(tr, 34, TEXT, bold=True)

    _accent_bar(slide)

    # Lead line (presenter / one-liner)
    if lead:
        lbox, ltf = _textbox(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6))
        lp = ltf.paragraphs[0]
        lr = lp.add_run(); lr.text = lead
        _set(lr, 15, MUTED, italic=True)

    # Bullets
    bbox, btf = _textbox(slide, Inches(0.7), Inches(2.5), Inches(12.0), Inches(4.3))
    btf.word_wrap = True
    first = True
    for text, level in bullets:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        bullet = "▸ " if level == 0 else "– "
        r = p.add_run(); r.text = bullet + text
        _set(r, 20 if level == 0 else 17, TEXT if level == 0 else MUTED, bold=False)

    _footer(slide, idx, total)
    return slide


def closing_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    block = slide.shapes.add_shape(1, 0, Inches(2.5), Inches(0.35), Inches(2.4))
    _fill(block, ACCENT)
    block.shadow.inherit = False

    box, tf = _textbox(slide, Inches(0.9), Inches(2.4), Inches(11.6), Inches(3.0))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Thank you"
    _set(r, 48, TEXT, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = ("From a multilingual pipeline to an interactive globe to "
               "natural-language Q&A — Trendy Topic answers what the world asks AI.")
    _set(r2, 18, MUTED, italic=True)
    p3 = tf.add_paragraph()
    r3 = p3.add_run(); r3.text = "Shocka & Niciah"
    _set(r3, 18, TEXT, bold=True)
    p4 = tf.add_paragraph()
    r4 = p4.add_run(); r4.text = "github.com/niciahrymer-hillian/trendy-topic"
    _set(r4, 16, ACCENT)
    _footer(slide, idx, total)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # (kicker, title, lead, [(bullet, level), ...])
    slides = [
        ("Problem", "The Problem",
         "Lead: Shocka",
         [("Millions talk to AI daily, but those conversations are a black box.", 0),
          ("We don't know what people actually ask — or how it differs by country and language.", 1),
          ("Raw logs are huge, messy, multilingual, and privacy-sensitive — unusable as-is.", 0),
          ("Research question: what does the world ask AI, and how do interests, tone, and language differ across countries?", 0)]),

        ("Solution", "Our Solution",
         "Lead: Niciah",
         [("Trendy Topic: an interactive platform turning raw AI conversations into explorable, privacy-safe insight.", 0),
          ("Combines data engineering + NLP + LLM analytics + geospatial visualization in one dashboard.", 0),
          ("Interactive globe, country & language analysis, sentiment, Curiosity Index, Q&A, translation, and AI voice.", 0)]),

        ("Data", "The Data",
         "Lead: Shocka",
         [("Source: WildChat real-world conversation corpus (Hugging Face, allenai/WildChat).", 0),
          ("Working sample: 12,000 conversations across 8 countries.", 0),
          ("Brazil, Canada, China, France, Great Britain, Japan, Russia, USA.", 1),
          ("Native languages preserved (EN, PT, FR, ZH, JA, RU) — real signal, not English-only.", 0),
          ("Formats supported: CSV, JSON, Parquet (streamed).", 0)]),

        ("Architecture", "System Architecture",
         "Lead: Niciah",
         [("Frontend: React + Vite + TypeScript with ECharts (port 5173).", 0),
          ("Backend: FastAPI Python service exposing analytics endpoints (port 8000).", 0),
          ("Processing: Pandas cleaning/enrichment, shared multilingual classifier, VADER sentiment, optional PostgreSQL.", 0),
          ("AI services: LLM extraction, embeddings for similarity/clustering, translation, ElevenLabs voice.", 0),
          ("One-click run: VS Code 'Start Trendy Topic Stack' boots API + UI together.", 0)]),

        ("Pipeline", "Data Pipeline",
         "Lead: Shocka",
         [("Ingest: stream WildChat rows by country into a common schema.", 0),
          ("Clean: normalize text, drop dupes/bad timestamps, mask PII (emails, tokens, secret-like blobs).", 0),
          ("Enrich: language detection → multilingual topic classification → sentiment.", 0),
          ("Classify smartly: boundary-aware keywords across 7 languages (basketball → Sports, trains → Transportation).", 0),
          ("Serve: tidy tables behind FastAPI; validated at 0 errors / 0 dupes / 0 bad timestamps.", 0)]),

        ("Dashboard", "Interactive Dashboard",
         "Lead: Niciah",
         [("Interactive Globe: spin and click a country to fly in and load its data.", 0),
          ("Global Overview: totals, redaction %, ranked topics, topic hierarchy tree.", 0),
          ("Compare Countries: two nations side-by-side on dual globes.", 0),
          ("Explore (tabbed): Topics trends · Languages heatmap · Sentiment split.", 0)]),

        ("Insights", "Key Insights",
         "Lead: Shocka",
         [("Global Curiosity Index: the most-asked questions across the corpus, ranked — our signature finding.", 0),
          ("Question Heatmap: topic intensity by country at a glance.", 0),
          ("Cross-cultural signal: Brazil over-indexes on Sports; Japan leans into Translation & Language.", 0),
          ("AI Insights: LLM-extracted trends & story angles; embeddings cluster countries by what they ask.", 0),
          ("Ask the Dataset: plain-English questions answered against live data.", 0)]),

        ("Ethics", "Ethics & Privacy",
         "Lead: Niciah",
         [("Privacy by design: raw sensitive data is never displayed; all views are aggregated.", 0),
          ("PII masking runs before analysis; a redaction % is shown transparently.", 0),
          ("Translation operates on safe summaries, never raw personal prompts.", 0),
          ("Multilingual fairness: native languages preserved so no community is misrepresented.", 0)]),

        ("Portfolio Value", "Portfolio Value",
         "Lead: Shocka & Niciah",
         [("Engineering breadth: end-to-end data eng + NLP + LLM + full-stack dashboard.", 0),
          ("Production habits: 235 passing tests, secret-scanning, type-checked frontend, one-click run.", 0),
          ("Real-world relevance: serves researchers, consultancies, and product teams.", 0),
          ("Storytelling: turns 3M+ potential conversations into an accessible, privacy-safe narrative.", 0)]),
    ]

    total = len(slides) + 2  # + title + close

    title_slide(prs)
    for i, (kicker, title, lead, bullets) in enumerate(slides, start=2):
        content_slide(prs, i, total, kicker, title, lead, bullets)
    closing_slide(prs, total, total)

    out = Path(__file__).resolve().parent.parent / "docs" / "Trendy_Topic_Capstone.pptx"
    prs.save(out)
    print(f"Wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
